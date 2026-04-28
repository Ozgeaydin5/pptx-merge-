from flask import Flask, request, jsonify, send_from_directory
import os, time, sys
from pptx import Presentation
from pptx.util import Inches
import fitz  # PyMuPDF

app = Flask(__name__)

BASE_DIR = "/home/ozge/mysite"
STATIC_DIR = os.path.join(BASE_DIR, "static")
PROJECT_DIR = os.path.join(BASE_DIR, "project")

# Kütüphane yolunu ekle
sys.path.append(os.path.expanduser("~/.local/lib/python3.10/site-packages"))

os.makedirs(STATIC_DIR, exist_ok=True)
os.makedirs(PROJECT_DIR, exist_ok=True)

@app.route("/merge", methods=["POST"])
def merge_to_pptx():
    try:
        from PyPDF2 import PdfMerger
        files = request.files.getlist("files")
        if not files:
            return jsonify({"error": "Dosya bulunamadı"}), 400

        # 1. ADIM: PDF'leri Birleştir
        merger = PdfMerger()
        for f in files:
            if f.filename:
                path = os.path.join(PROJECT_DIR, f.filename)
                f.save(path)
                merger.append(path)

        temp_pdf = os.path.join(PROJECT_DIR, "temp_merged.pdf")
        with open(temp_pdf, "wb") as f:
            merger.write(f)
        merger.close()

        # 2. ADIM: Birleşen PDF'i PPTX'e Çevir
        prs = Presentation()
        # Standart 16:9 ekran boyutları
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        doc = fitz.open(temp_pdf)
        for page in doc:
            # Sayfayı resme çevir (kalite için 2x zoom)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_path = os.path.join(PROJECT_DIR, f"temp_page_{page.number}.png")
            pix.save(img_path)

            # Slayt ekle ve resmi tam ekran yerleştir
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # 6: Boş slayt düzeni
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            os.remove(img_path) # Geçici resmi sil
        doc.close()

        # 3. ADIM: Kaydet ve Temizle
        dosya_adi = f"birlesik_sunum_{int(time.time())}.pptx"
        output_path = os.path.join(STATIC_DIR, dosya_adi)
        prs.save(output_path)
        
        os.remove(temp_pdf) # Geçici PDF'i sil

        return jsonify({
            "izle": f"https://ozge.pythonanywhere.com/static/{dosya_adi}",
            "indir": f"https://ozge.pythonanywhere.com/indir/{dosya_adi}",
            "deneme": f"https://ozge.pythonanywhere.com/indir/{dosya_adi}"
        }), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/indir/<filename>")
def download_file(filename):
    return send_from_directory(STATIC_DIR, filename, as_attachment=True)

if __name__ == "__main__":
    app.run()