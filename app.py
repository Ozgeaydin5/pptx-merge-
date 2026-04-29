from flask import Flask, request, jsonify
import requests
from pptx import Presentation
import io
import os

app = Flask(__name__)

def copy_slides(src, dest):
    """Önceki uygulamadaki slayt kopyalama mantığını koruyoruz."""
    for slide in src.slides:
        # Slayt düzenini (layout) hedef sunuma uydurarak ekle
        slide_layout = dest.slide_layouts[src.slide_layouts.index(slide.slide_layout)]
        new_slide = dest.slides.add_slide(slide_layout)
        
        # Slayt üzerindeki her bir objeyi (şekil, metin vb.) kopyala
        for shape in slide.shapes:
            if shape.has_text_frame:
                new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                new_shape.text = shape.text
            # Resim veya grafik varsa burada genişletilebilir

@app.route('/merge', methods=['POST'])
def merge_pptx():
    try:
        data = request.get_json()
        # NocoBase'den gelen "urls" listesini yakala
        urls = data.get('urls', [])

        if not urls or not isinstance(urls, list):
            return jsonify({"error": "Dosya listesi boş veya hatalı formatta"}), 400

        # Ana sunumu oluştur (İlk sunumu temel alabilirsin veya boş başlatabilirsin)
        merged_pptx = Presentation()
        # İlk boş slaytı sil (Yeni sunum oluştururken default gelen)
        if len(merged_pptx.slides) > 0:
            xml_slides = merged_pptx.slides._slim_elements
            del xml_slides[0]

        for url in urls:
            response = requests.get(url)
            if response.status_code == 200:
                pptx_file = io.BytesIO(response.content)
                current_pptx = Presentation(pptx_file)
                # Kopyalama fonksiyonunu çağır
                copy_slides(current_pptx, merged_pptx)
            else:
                print(f"Hata: {url} indirilemedi.")

        # Birleşmiş dosyayı belleğe kaydet
        output = io.BytesIO()
        merged_pptx.save(output)
        output.seek(0)

        # --- BURASI ÖNEMLİ ---
        # Dosyayı NocoBase'e geri göndermek için bir yere upload etmen gerekir.
        # Eğer bir storage kullanmıyorsan, NocoBase bu dosyayı URL olarak bekler.
        # Geçici olarak bir dosya servisi veya kendi sunucunu kullanabilirsin.
        # Örnek dönüş:
        return jsonify({
            "output_url": "BIRLESEN_DOSYA_LINKI_BURAYA_GELECEK"
        }), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5000))
    app.run(host='0.0.0.0', port=port)
